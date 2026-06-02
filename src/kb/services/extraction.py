"""File-type-specific text extraction, preserving page boundaries.

Each extractor returns a list of (page_number, text) tuples. For flat formats
(CSV, XLSX) each row or sheet maps to a "page". For paginated formats (PDF,
PPTX, DOCX) the page/slide/section index is preserved.

Direct text extraction is attempted first; OCR is used as fallback when a
page yields little text but contains images (scanned content). For PDFs,
tables detected by pymupdf are rendered as pipe-grids so the LLM can still
see column/row relationships rather than a flattened token stream.
"""

from __future__ import annotations

import csv
import io
import logging
from pathlib import Path
from typing import Any

log = logging.getLogger("kb.extraction")

PageText = tuple[int, str]  # (page_number, text)


class ScannedPdfError(ValueError):
    """Raised when a PDF has no readable text and looks like a scan/image.

    Distinct from a generic extraction failure so the import pipeline can tell
    the user the file needs OCR rather than reporting an opaque error.
    """


def _try_import(module: str) -> Any:
    """Lazy-import an optional dependency, returning None if missing."""
    try:
        import importlib
        return importlib.import_module(module)
    except ImportError:
        return None


# ── PDF extraction ───────────────────────────────────────────────────────────

def extract_pdf(
    path: Path,
    *,
    ocr_enabled: bool = True,
    ocr_lang: str = "ch",
    ocr_min_confidence: float = 0.5,
) -> list[PageText]:
    fitz = _try_import("fitz")
    if fitz is None:
        raise ImportError("pymupdf is required for PDF extraction: pip install pymupdf")

    pages: list[PageText] = []
    blank_pages = 0          # pages that yielded no usable text at all
    image_only_pages = 0     # image-dominated pages we couldn't read as text
    doc = fitz.open(str(path))
    try:
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text("text") or ""

            # Augment with table grids so vertical/horizontal tables retain
            # column-row structure when fed to the segmentation LLM. Tables
            # already appear as a flat token stream in `text`; we append a
            # structured rendering as well — the LLM gets both views.
            table_md = _extract_pdf_tables(page)
            if table_md:
                text = f"{text}\n\n{table_md}" if text.strip() else table_md

            wanted_ocr = _should_use_ocr(text, page, ocr_enabled)
            if wanted_ocr:
                from kb.services.ocr import ocr_page_image
                pix = page.get_pixmap(dpi=300)
                img_bytes = pix.tobytes("png")
                try:
                    ocr_text = ocr_page_image(
                        img_bytes, lang=ocr_lang, min_confidence=ocr_min_confidence
                    )
                except Exception as exc:  # noqa: BLE001 — OCR is best-effort
                    log.warning("OCR failed on page %d: %s", page_num + 1, exc)
                    ocr_text = ""
                # Prefer OCR only if it's *meaningfully* longer AND looks like
                # real text (avoids replacing real text with OCR garbage).
                if (
                    len(ocr_text.strip()) > len(text.strip()) * 1.2
                    and _looks_like_text(ocr_text)
                ):
                    text = ocr_text

            if text.strip():
                pages.append((page_num + 1, _clean_extracted_text(text)))
            else:
                blank_pages += 1
                # No text AND the page carries images → it's a scan we couldn't
                # read (OCR off, unavailable, or low-quality). Track separately
                # so the caller can give the user an actionable message instead
                # of a bare "No text extracted".
                if _page_has_images(page):
                    image_only_pages += 1
    finally:
        doc.close()

    if blank_pages:
        log.info(
            "PDF %s: %d/%d page(s) yielded no text (%d image-only)",
            path.name, blank_pages, len(doc), image_only_pages,
        )
    # Surface a likely-scanned document so the pipeline can advise on OCR rather
    # than silently returning fewer documents than the source contains.
    if not pages and image_only_pages:
        raise ScannedPdfError(
            f"No readable text in {path.name} — it appears to be a scanned/image PDF "
            f"({image_only_pages} image page(s))."
        )
    return pages


def _extract_pdf_tables(page: Any) -> str:
    """Render every detected table on the page as a pipe-grid.

    pymupdf's table finder handles both horizontal (header-on-top) and
    vertical (header-on-left) layouts — orientation is inferred from cell
    geometry. We don't try to second-guess it: whatever pymupdf returns,
    we render as-is so the downstream LLM gets a 2D view instead of a
    flattened token soup.

    Returns an empty string on any error — table augmentation is best-effort.
    """
    try:
        finder = page.find_tables()
        tables = getattr(finder, "tables", finder) or []
    except Exception as exc:  # noqa: BLE001 — best-effort
        log.debug("find_tables failed: %s", exc)
        return ""

    parts: list[str] = []
    for t_idx, table in enumerate(tables, start=1):
        try:
            rows = table.extract()
        except Exception:  # noqa: BLE001
            continue
        # Drop fully-empty rows; collapse Nones to "".
        clean_rows = [
            [(cell or "").strip().replace("\n", " ").replace("|", "/") for cell in row]
            for row in rows
            if any((c or "").strip() for c in row)
        ]
        if not clean_rows:
            continue
        rendered = "\n".join("| " + " | ".join(row) + " |" for row in clean_rows)
        parts.append(f"[Table {t_idx}]\n{rendered}")
    return "\n\n".join(parts)


def _looks_like_text(s: str) -> bool:
    """Heuristic: is this string mostly real characters vs. OCR garbage?"""
    if not s:
        return False
    printable = sum(1 for c in s if c.isprintable() or c in "\n\r\t")
    return (printable / len(s)) >= 0.85


# Common PDF extraction artifacts: NULs, soft hyphens, BOMs, form feeds.
_PDF_NOISE_CHARS = str.maketrans({
    "\x00": "",   # NUL
    "\xad": "",   # soft hyphen — invisible, breaks search
    "﻿": "", # BOM
    "\x0c": "\n", # form feed → newline
    "\x0b": "\n", # vertical tab → newline
})


def _clean_extracted_text(text: str) -> str:
    """Strip control-character noise that PDF extractors leak.

    Keeps tab/newline/CR; removes NUL and other C0 controls; normalizes
    Windows/Mac line endings; collapses runs of >2 blank lines.
    """
    text = text.translate(_PDF_NOISE_CHARS)
    # Strip remaining C0 controls except \t \n \r
    text = "".join(
        c for c in text
        if c >= " " or c in "\t\n\r"
    )
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    # Collapse 3+ consecutive newlines to 2
    while "\n\n\n" in text:
        text = text.replace("\n\n\n", "\n\n")
    return text


def _page_has_images(page: Any) -> bool:
    """Whether the page carries any embedded raster images."""
    try:
        return len(page.get_images(full=True)) > 0
    except Exception:  # noqa: BLE001 — best-effort
        return False


def _image_coverage(page: Any) -> float:
    """Fraction of the page area covered by embedded images (0.0–1.0).

    A scanned page is typically one full-page image; a digital page with a
    small logo is not. Used to OCR image-dominated pages even when the text
    layer leaked some junk characters.
    """
    try:
        page_area = abs(page.rect.width * page.rect.height)
        if page_area <= 0:
            return 0.0
        img_area = 0.0
        for info in page.get_image_info():
            bbox = info.get("bbox")
            if bbox:
                x0, y0, x1, y1 = bbox
                img_area += abs((x1 - x0) * (y1 - y0))
        return min(img_area / float(page_area), 1.0)
    except Exception:  # noqa: BLE001 — best-effort
        return 0.0


# Below this many extracted chars a page is treated as "text-sparse" — if it's
# also image-dominated we OCR it. Higher than the old 50 so scanned pages that
# leak a header/footer line still get OCR'd.
_OCR_TEXT_THRESHOLD = 120


def _should_use_ocr(text: str, page: Any, ocr_enabled: bool) -> bool:
    if not ocr_enabled:
        return False
    text_chars = len(text.strip())
    has_images = _page_has_images(page)
    # Sparse text on a page that has images → likely a scan with a thin text
    # layer. Use coverage so a page that's mostly one big image is OCR'd even
    # if a few stray characters were extracted.
    if has_images and (text_chars < _OCR_TEXT_THRESHOLD or _image_coverage(page) >= 0.5):
        return True
    # Text present but mostly non-printable → garbled extraction, re-OCR.
    if text_chars > 0:
        printable_ratio = sum(1 for c in text if c.isprintable() or c in "\n\r\t") / len(text)
        if printable_ratio < 0.7:
            return True
    return False


# ── PPTX extraction ─────────────────────────────────────────────────────────

def extract_pptx(path: Path) -> list[PageText]:
    pptx_mod = _try_import("pptx")
    if pptx_mod is None:
        raise ImportError("python-pptx is required for PPTX extraction: pip install python-pptx")

    from pptx import Presentation

    pages: list[PageText] = []
    prs = Presentation(str(path))
    for slide_num, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = paragraph.text.strip()
                    if line:
                        texts.append(line)
            if shape.has_table:
                table = shape.table
                rendered = _render_office_table(
                    [[cell.text for cell in row.cells] for row in table.rows]
                )
                if rendered:
                    texts.append(rendered)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                texts.append(f"[Notes] {notes}")
        if texts:
            pages.append((slide_num, "\n".join(texts)))
    return pages


# ── DOCX extraction ─────────────────────────────────────────────────────────

def extract_docx(path: Path) -> list[PageText]:
    docx_mod = _try_import("docx")
    if docx_mod is None:
        raise ImportError("python-docx is required for DOCX extraction: pip install python-docx")

    from docx import Document
    from docx.oxml.ns import qn
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    doc = Document(str(path))

    # Walk body in document order so paragraphs and tables stay interleaved
    # the way a human reading the doc sees them. The previous "all paragraphs
    # then all tables" pass destroyed locality between an entry written in
    # prose and a following one-row table that named its Project / Equipment.
    p_tag = qn("w:p")
    tbl_tag = qn("w:tbl")
    body = doc.element.body
    texts: list[str] = []
    for child in body.iterchildren():
        if child.tag == p_tag:
            t = Paragraph(child, body).text.strip()
            if t:
                texts.append(t)
        elif child.tag == tbl_tag:
            tbl = Table(child, body)
            rendered = _render_office_table(
                [[cell.text for cell in row.cells] for row in tbl.rows]
            )
            if rendered:
                texts.append(rendered)
    full_text = "\n".join(texts)
    if not full_text.strip():
        return []
    return [(1, full_text)]


def _render_office_table(rows: list[list[str]]) -> str:
    """Render a DOCX/PPTX table as a pipe-grid, preserving cell layout."""
    clean = [
        [(cell or "").strip().replace("\n", " ").replace("|", "/") for cell in row]
        for row in rows
        if any((c or "").strip() for c in row)
    ]
    if not clean:
        return ""
    return "\n".join("| " + " | ".join(row) + " |" for row in clean)


# ── XLSX / XLS extraction ───────────────────────────────────────────────────

def extract_xlsx(path: Path) -> list[PageText]:
    openpyxl = _try_import("openpyxl")
    if openpyxl is None:
        raise ImportError("openpyxl is required for XLSX extraction: pip install openpyxl")

    from openpyxl import load_workbook

    # read_only mode keeps the workbook file open until close(); use try/finally
    # so an error mid-iteration still releases the handle.
    wb = load_workbook(str(path), read_only=True, data_only=True)
    pages: list[PageText] = []
    try:
        for sheet_idx, sheet_name in enumerate(wb.sheetnames, start=1):
            ws = wb[sheet_name]
            rows: list[str] = []
            # Tag with sheet name so the LLM can disambiguate when multiple sheets
            # carry similar structure (e.g. "Alarms_EN" vs "Alarms_ZH").
            rows.append(f"[Sheet: {sheet_name}]")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c).strip() if c is not None else "" for c in row]
                line = "| " + " | ".join(cells) + " |"
                if any(c for c in cells):
                    rows.append(line)
            if len(rows) > 1:  # more than just the sheet header
                pages.append((sheet_idx, "\n".join(rows)))
    finally:
        wb.close()
    return pages


# ── CSV extraction ───────────────────────────────────────────────────────────

def extract_csv(path: Path) -> list[PageText]:
    for encoding in ("utf-8-sig", "utf-8", "gb18030", "latin-1"):
        try:
            text = path.read_text(encoding=encoding)
            break
        except (UnicodeDecodeError, ValueError):
            continue
    else:
        return []

    reader = csv.reader(io.StringIO(text))
    rows: list[str] = []
    for row in reader:
        line = "\t".join(c.strip() for c in row)
        if line.replace("\t", "").strip():
            rows.append(line)
    if not rows:
        return []
    return [(1, "\n".join(rows))]


# ── Dispatcher ───────────────────────────────────────────────────────────────

EXTRACTORS = {
    "pdf": extract_pdf,
    "pptx": extract_pptx,
    "docx": extract_docx,
    "xlsx": extract_xlsx,
    "xls": extract_xlsx,
    "csv": extract_csv,
}


def extract_file(
    path: Path,
    *,
    ocr_enabled: bool = True,
    ocr_lang: str = "ch",
    ocr_min_confidence: float = 0.5,
) -> list[PageText]:
    """Extract text from a file, returning (page_number, text) pairs.

    Raises ImportError if the required library is not installed.
    Raises ValueError for unsupported file types.
    Raises ScannedPdfError if a PDF has no readable text and looks like a scan.
    """
    suffix = path.suffix.lower().lstrip(".")
    extractor = EXTRACTORS.get(suffix)
    if extractor is None:
        raise ValueError(f"Unsupported file type: {suffix}")

    if suffix == "pdf":
        return extract_pdf(
            path,
            ocr_enabled=ocr_enabled,
            ocr_lang=ocr_lang,
            ocr_min_confidence=ocr_min_confidence,
        )
    return extractor(path)
