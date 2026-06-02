"""Session-scoped fixture-file generators.

Rather than committing large binary PPTX/PDF templates to the repo, these
fixtures build minimal but content-complete files on-the-fly from the seed
CSVs.  Each generator runs once per pytest session and saves the result to a
pytest-managed tmp directory; the files are cleaned up automatically after the
session ends.

Generated files are deliberately plain/minimal (no images, no themes) so they
exercise text-extraction code paths without the per-run cost of rendering.
"""
from __future__ import annotations

import csv
import re
from collections import Counter, defaultdict
from itertools import groupby
from pathlib import Path
from typing import Any

import pytest

CONFIG_DIR = Path(__file__).parent.parent.parent / "config"

# ---------------------------------------------------------------------------
# Shared CSV helpers
# ---------------------------------------------------------------------------

def _read_csv(path: Path) -> list[dict[str, str]]:
    for enc in ("utf-8-sig", "utf-8", "gb18030"):
        try:
            with open(path, encoding=enc) as f:
                return list(csv.DictReader(f))
        except (UnicodeDecodeError, ValueError):
            continue
    return []


@pytest.fixture(scope="session")
def alarm_csv_data() -> list[dict[str, str]]:
    return _read_csv(CONFIG_DIR / "机台报警_header.csv")


@pytest.fixture(scope="session")
def setup_csv_data() -> list[dict[str, str]]:
    return _read_csv(CONFIG_DIR / "机台setup_header.csv")


@pytest.fixture(scope="session")
def experience_csv_data() -> list[dict[str, str]]:
    return _read_csv(CONFIG_DIR / "设备经验_header.csv")


# ---------------------------------------------------------------------------
# Alarm PPTX generator
# ---------------------------------------------------------------------------

@pytest.fixture(scope="session")
def alarm_pptx_path(tmp_path_factory, alarm_csv_data):
    """Minimal PPTX generated from the alarm CSV.

    Slide structure:
    - Slide 1: cover with the three severity-level labels (not in CSV; come
      from the structural template pages in the original PPTX).
    - Slide 2: equipment-type and project-name roster.
    - One slide per alarm row: code, titles, equipment, project, 内容,
      解除流程, 注意事项.
    """
    pytest.importorskip("pptx", reason="python-pptx not installed")
    from pptx import Presentation

    prs = Presentation()
    blank = prs.slide_layouts[6]

    def _slide(text: str) -> None:
        slide = prs.slides.add_slide(blank)
        box = slide.shapes.add_textbox(0, 0, prs.slide_width, prs.slide_height)
        box.text_frame.word_wrap = True
        box.text_frame.text = text

    # Cover — structural severity vocabulary lives here, not in per-alarm rows.
    _slide(
        "机台报警 · ALARM COMPENDIUM\n"
        "严重程度三级编排 · 含报警原理、解除流程与安全警示\n"
        "\n"
        "01 安全联锁 SAFETY INTERLOCK\n"
        "硬联锁触发，涉及人身安全或不可逆机械损伤。\n"
        "\n"
        "02 故障停机 FAULT STOP\n"
        "设备检测到故障并保护性停机。须按解除流程排查根因。\n"
        "\n"
        "03 预警提示 WARNING\n"
        "预测性、合规或保养类提示。"
    )

    # Roster slide
    equips = " · ".join(sorted({r["机台"] for r in alarm_csv_data}))
    projs = " · ".join(sorted({r["项目"] for r in alarm_csv_data}))
    _slide(f"机台类型: {equips}\n产线: {projs}")

    # One slide per alarm — include every field the tests assert against.
    for row in alarm_csv_data:
        _slide(
            f"{row['代码']}  {row['英文标题']}  {row['中文标题']}\n"
            f"机台: {row['机台']}  产线: {row['项目']}\n"
            f"内容: {row['内容'][:300]}\n"
            f"解除流程:\n{row['解除流程'][:300]}\n"
            f"注意事项: {row['注意事项'][:150]}"
        )

    path = tmp_path_factory.mktemp("alarm") / "alarm.pptx"
    prs.save(str(path))
    return path


# ---------------------------------------------------------------------------
# Setup PPTX generator
# ---------------------------------------------------------------------------

@pytest.fixture(scope="session")
def setup_pptx_path(tmp_path_factory, setup_csv_data):
    """Minimal PPTX generated from the setup CSV.

    Slide structure:
    - Slide 1: cover listing each equipment type with its station count
      (e.g. "Aligner 9 工站") — the test hard-checks 9 and 7.
    - Per-equipment section slide with SPEC/PROCEDURE/TOOLS headers.
    - Per-station table slide with one row per column.
    """
    pytest.importorskip("pptx", reason="python-pptx not installed")
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    blank = prs.slide_layouts[6]

    def _text_slide(text: str) -> None:
        slide = prs.slides.add_slide(blank)
        box = slide.shapes.add_textbox(0, 0, prs.slide_width, prs.slide_height)
        box.text_frame.word_wrap = True
        box.text_frame.text = text

    def _table_slide(header: str, rows: list[list[str]]) -> None:
        slide = prs.slides.add_slide(blank)
        hbox = slide.shapes.add_textbox(0, 0, prs.slide_width, Inches(0.5))
        hbox.text_frame.text = header
        if rows:
            tbl = slide.shapes.add_table(
                len(rows), 2, 0, Inches(0.55), prs.slide_width, Inches(3.5)
            ).table
            for r_i, (label, value) in enumerate(rows):
                tbl.cell(r_i, 0).text = label
                tbl.cell(r_i, 1).text = value

    # Cover: station counts per equipment type.
    counts = Counter(r["设备"] for r in setup_csv_data)
    cover_lines = ["机台 Setup 调试说明 · SETUP GUIDE"]
    for equip in sorted(counts):
        cover_lines.append(f"{equip} {counts[equip]} 工站")
    _text_slide("\n".join(cover_lines))

    # Per-equipment section + per-station slides.
    for equip, grp in groupby(
        sorted(setup_csv_data, key=lambda r: r["设备"]),
        key=lambda r: r["设备"],
    ):
        _text_slide(
            f"{equip}\n"
            "SPEC / 规格  ·  PROCEDURE / 调试步骤  ·  TOOLS / 调试工具"
        )
        for row in grp:
            _table_slide(
                f"{equip} · {row['工站/部件/站位']}",
                [
                    ["SPEC 规格",          row["规格/要求"][:120]],
                    ["PROCEDURE 调试步骤", row["调试步骤"][:120]],
                    ["TOOLS 调试工具",     row["调试工具"][:100]],
                ],
            )

    path = tmp_path_factory.mktemp("setup") / "setup.pptx"
    prs.save(str(path))
    return path


# ---------------------------------------------------------------------------
# Experience PDF generator
# ---------------------------------------------------------------------------

# Category metadata — mirrors the eight-category taxonomy in the original doc.
_CATEGORIES: list[tuple[str, str, str]] = [
    ("A", "电气",  "ELECTRICAL & SIGNAL INTEGRITY"),
    ("B", "机械",  "MECHANICAL & WEAR"),
    ("C", "寿命",  "END-OF-LIFE & AGING"),
    ("D", "热管理","THERMAL MANAGEMENT"),
    ("E", "污染",  "CONTAMINATION & PARTICLES"),
    ("F", "选型",  "SELECTION & DESIGN"),
    ("G", "厂务",  "FACILITY & ENVIRONMENT"),
    ("H", "管理",  "MANAGEMENT & PROCESS"),
]

# Keyword patterns used to assign each CSV row to a category.
# Order matters — first match wins.  D (thermal) must precede B (mechanical)
# because thermal root-causes often contain the word 机械 (e.g. "热膨胀导致
# 相机机械结构漂移"), which would otherwise be captured by B first.
_CAT_PATTERNS: list[tuple[str, str]] = [
    ("D", r"热膨胀|散热不足|温升|热管理|SSR.*散热|散热.*SSR|硅脂|热阻"),
    ("A", r"电气|接地|EMI|ESD|网络|信号|电源|漂移噪|ADC|dropout|干扰.*传感|传感.*干扰"),
    ("B", r"机械|磨损|松动|同轴|润滑|紧固|螺钉|backlash|挡片|门钩|密封.*变形"),
    ("C", r"寿命末期|老化|轴承|BPFO|继电器寿|加热膜.*老|光耦.*寿|热电偶老"),
    ("E", r"污染|碎屑|颗粒|油气|油雾|雾化|尘|磁条.*金属"),
    ("F", r"选型|材质.*不兼容|不兼容.*材质|拓扑设计|设计.*缺陷|供应商批次"),
    ("G", r"厂务|气源|供电|环境光|地基|电压偏低|空调|气流|外部环境"),
    ("H", r"管理|流程|SOP|未受控|来料|PM 规范|扭矩|权限|变更|停机.*SOP"),
]


def _assign_category(root_cause: str) -> str:
    for letter, pattern in _CAT_PATTERNS:
        if re.search(pattern, root_cause):
            return letter
    return "H"  # default: management/process


def _pdf_add_page(doc: Any, lines: list[str]) -> None:
    """Append one A4 page and render ``lines`` using the built-in china-s font.

    ``china-s`` (GB2312 Simplified Chinese) is one of the fonts bundled with
    pymupdf and covers the full CJK range used by the seed CSVs — no external
    font file required.  Lines that would fall below the page bottom are
    silently dropped (the tests do not assert a specific page layout).
    """
    import fitz

    page = doc.new_page(width=595, height=842)
    y = 48.0
    for line in lines:
        if not line:
            y += 8.0
            continue
        if y > 820:
            break
        page.insert_text((40, y), line, fontname="china-s", fontsize=10)
        y += 15.0


@pytest.fixture(scope="session")
def experience_pdf_path(tmp_path_factory, experience_csv_data):
    """Minimal text-only PDF generated from the experience CSV.

    Page structure:
    - Page 1: cover (equipment types, production lines, case counts).
    - Page 2: root-cause taxonomy A–H.
    - Per category: summary page (all cases in tabular form) +
      deep-dive page (one representative case with full field labels and
      a PPT file reference).
    - Bulk-data pages: all 50 rows with full 问题 / 根因 / 纠正步骤 text,
      ensuring complete coverage regardless of category assignment.
    """
    fitz_mod = pytest.importorskip("fitz", reason="pymupdf not installed")
    import fitz

    doc = fitz.open()

    # Page 1: cover
    equips = " ".join(sorted({r["机台"] for r in experience_csv_data}))
    projs = " ".join(
        sorted({r.get("项目", "").replace("所有项目", "ALL") for r in experience_csv_data})
    )
    _pdf_add_page(doc, [
        "设备工程 · EQUIPMENT ENGINEERING",
        "设备经验知识库 / KNOWLEDGE BASE",
        "50 案例  8 机台类型  7 产线  8 根因类别 ROOT-CAUSE",
        "问题 → 失败描述 → 失败分析 → 根因 → 纠正措施 → PM 固化",
        f"机台: {equips}",
        f"产线: {projs}",
    ])

    # Page 2: full taxonomy A–H
    tax_lines = ["根因类别 / TAXONOMY  八大根因类别"]
    for letter, zh, en in _CATEGORIES:
        tax_lines.append(f"  {letter} · {zh}  {en}")
    _pdf_add_page(doc, tax_lines)

    # Assign each CSV row to a category.
    cat_rows: dict[str, list[dict]] = defaultdict(list)
    for row in experience_csv_data:
        cat_rows[_assign_category(row.get("根因", ""))].append(row)

    # Per-category: summary page + one deep-dive page.
    for letter, zh, en in _CATEGORIES:
        rows = cat_rows.get(letter, [])

        # Each column header goes on its own line so none are truncated by the
        # page-width limit — the tests check for each string independently.
        summary = [
            f"{letter} · {zh} / {en}",
            "问题 PROBLEM",
            "根因 ROOT CAUSE",
            "关键纠正 KEY ACTION",
        ]
        for row in rows:
            proj = row.get("项目", "").replace("所有项目", "ALL")
            step1 = row.get("纠正步骤", "").split("；")[0].lstrip("1) ").strip()[:50]
            summary.append(
                f"  {proj} {row.get('机台','')}  {row.get('问题','')}  "
                f"{row.get('根因','')[:50]}  → {step1}"
            )
        _pdf_add_page(doc, summary)

        if rows:
            row = rows[0]
            pptx_ref = row.get("PPT文件", f"exp_{letter}.pptx")
            proj = row.get("项目", "").replace("所有项目", "ALL LINES")
            _pdf_add_page(doc, [
                f"{letter} · {zh} / {en}",
                "案例精解 / D E E P  D I V E",
                f"{row.get('机台', '')} {row.get('问题', '')}",
                f"产线 LINE  {proj}",
                f"机台 EQUIPMENT  {row.get('机台', '')}",
                f"类别 CATEGORY  {letter} · {en}",
                f"参考 REF {pptx_ref} · p.{row.get('PPT页面', '')}",
                "",
                "失败描述 · SYMPTOM",
                row.get("失败描述", "")[:160],
                "",
                "失败分析 · ANALYSIS",
                row.get("失败分析", "")[:160],
                "",
                "根因 · ROOT CAUSE",
                row.get("根因", "")[:160],
                "",
                "纠正措施 · CORRECTIVE ACTION",
                row.get("纠正步骤", "")[:160],
            ])

    # Bulk-data pages: all 50 rows in full — guarantees every 问题 / 根因 /
    # 纠正步骤 value from the CSV appears in the extracted text regardless of
    # how the category-assignment heuristic distributed them above.
    chunk = 5
    for i in range(0, len(experience_csv_data), chunk):
        lines: list[str] = []
        for row in experience_csv_data[i : i + chunk]:
            lines += [
                f"【{row.get('问题', '')}】",
                f"  根因: {row.get('根因', '')}",
                f"  纠正: {row.get('纠正步骤', '').split(chr(10))[0][:100]}",
                "",
            ]
        _pdf_add_page(doc, lines)

    path = tmp_path_factory.mktemp("exp") / "experience.pdf"
    doc.save(str(path))
    doc.close()
    return path
